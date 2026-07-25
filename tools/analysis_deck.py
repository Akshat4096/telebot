"""
PPTX business analysis deck: sales trend, top items, stock health, GST
collected — built with python-pptx + matplotlib (real charts rendered to
PNG and embedded, not screenshots of some other tool).
"""
from __future__ import annotations

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "data" / "generated"
BRAND = RGBColor(0x2D, 0x34, 0x36)
ACCENT = RGBColor(0x00, 0x9E, 0x60)


def _chart_png(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.6), Inches(1.5))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BRAND
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x63, 0x63, 0x63)
    return slide


def _bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = BRAND

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(18)
    return slide


def _chart_slide(prs, title, fig, insight: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = BRAND

    png = _chart_png(fig)
    # Fit the rendered chart into a fixed box (whatever its aspect ratio turns
    # out to be — pie charts render nearly square, bar charts wide) instead of
    # forcing a fixed width, which previously let tall/square charts overflow
    # off the bottom of the slide and collide with the insight caption below.
    from PIL import Image
    png.seek(0)
    px_w, px_h = Image.open(png).size
    png.seek(0)
    max_w, max_h = Inches(8.4), Inches(5.15)
    scale = min(max_w / px_w, max_h / px_h)
    w, h = int(px_w * scale), int(px_h * scale)
    left = int((Inches(10) - w) / 2)
    top = Inches(1.2)
    slide.shapes.add_picture(png, left, top, width=w, height=h)
    if insight:
        note = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(8.4), Inches(0.7))
        note.text_frame.text = insight
        note.text_frame.paragraphs[0].font.size = Pt(13)
        note.text_frame.paragraphs[0].font.italic = True
        note.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x63, 0x63, 0x63)
    return slide


def build_deck(shop_name: str, date_from: str, date_to: str, summary: dict,
               low_stock: list, title: str | None = None) -> str:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs, title or f"{shop_name} — Sales Analysis",
                 f"{date_from} to {date_to}")

    # Slide: headline numbers
    _bullet_slide(prs, "Headline numbers", [
        f"Total sales: ₹{summary['total_sales']:.2f} across {summary['bill_count']} bills",
        f"GST collected: ₹{summary['total_gst']:.2f}  (CGST ₹{summary['total_cgst']:.2f} + "
        f"SGST ₹{summary['total_sgst']:.2f})",
        f"Average bill value: ₹{(summary['total_sales']/summary['bill_count']):.2f}"
        if summary['bill_count'] else "No finalized bills in this period",
        f"Payment split: " + ", ".join(f"{k.upper()} ₹{v:.2f}" for k, v in summary['by_payment_mode'].items())
        if summary['by_payment_mode'] else "No payments recorded",
    ])

    # Chart: payment mode split (pie)
    if summary["by_payment_mode"]:
        fig, ax = plt.subplots(figsize=(6, 4.2))
        modes = list(summary["by_payment_mode"].keys())
        values = list(summary["by_payment_mode"].values())
        ax.pie(values, labels=[m.upper() for m in modes], autopct="%1.0f%%",
               colors=["#009E60", "#2D3436", "#F5A623", "#5B8DEF"][:len(modes)])
        ax.set_title("Sales by payment mode")
        _chart_slide(prs, "Payment mode split", fig,
                     "Cash vs UPI vs Card vs Credit share of total revenue.")

    # Chart: top items (bar)
    if summary["top_items"]:
        fig, ax = plt.subplots(figsize=(7.5, 4.2))
        names = [t["name"][:18] for t in summary["top_items"]][::-1]
        revenue = [t["revenue"] for t in summary["top_items"]][::-1]
        ax.barh(names, revenue, color="#009E60")
        ax.set_xlabel("Revenue (₹)")
        ax.set_title("Top-selling items by revenue")
        _chart_slide(prs, "Top items", fig, f"Best seller: {summary['top_items'][0]['name']}.")

    # Chart: stock health
    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    if low_stock:
        names = [p["name"][:18] for p in low_stock][:10][::-1]
        levels = [p["quantity"] for p in low_stock][:10][::-1]
        reorder = [p["reorder_level"] for p in low_stock][:10][::-1]
        y = range(len(names))
        ax.barh(list(y), levels, color="#F5A623", label="Current stock")
        ax.barh(list(y), reorder, color="none", edgecolor="#2D3436", linewidth=1.2, label="Reorder level")
        ax.set_yticks(list(y))
        ax.set_yticklabels(names)
        ax.legend()
        ax.set_title("Items at/below reorder level")
    else:
        ax.text(0.5, 0.5, "Nothing below reorder level", ha="center", va="center", fontsize=14)
        ax.axis("off")
    _chart_slide(prs, "Stock health", fig,
                 f"{len(low_stock)} SKU(s) need reordering." if low_stock else "Stock levels are healthy.")

    path = OUTPUT_DIR / f"analysis_{date_from}_{date_to}.pptx"
    prs.save(str(path))
    return str(path)
